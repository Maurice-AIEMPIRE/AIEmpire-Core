"""
Office Processor - Word, Excel, PowerPoint
"""

from typing import Dict, Any
from .base import BaseProcessor
import logging

logger = logging.getLogger(__name__)


class OfficeProcessor(BaseProcessor):
    """Verarbeitet Office-Dateien (docx, xlsx, pptx)"""

    async def extract(self) -> Dict[str, Any]:
        """Extrahiere Text aus Office-Dateien"""
        file_type = self.file_path.suffix.lower()

        if file_type == ".docx":
            return await self._process_docx()
        elif file_type == ".xlsx":
            return await self._process_xlsx()
        elif file_type == ".pptx":
            return await self._process_pptx()
        else:
            return {"error": f"Office-Typ nicht unterstützt: {file_type}"}

    async def _process_docx(self) -> Dict[str, Any]:
        """Word-Datei"""
        try:
            from docx import Document
        except ImportError:
            return {"error": "python-docx nicht installiert"}

        try:
            doc = Document(self.file_path)
            paragraphs = [p.text for p in doc.paragraphs]
            tables = []
            for table in doc.tables:
                table_data = [
                    [cell.text for cell in row.cells] for row in table.rows
                ]
                tables.append(table_data)

            return {
                "content_type": "docx",
                "text": "\n".join(paragraphs),
                "tables": tables,
            }
        except Exception as e:
            return {
                "content_type": "docx",
                "error": str(e),
            }

    async def _process_xlsx(self) -> Dict[str, Any]:
        """Excel-Datei"""
        try:
            import openpyxl
        except ImportError:
            return {"error": "openpyxl nicht installiert"}

        try:
            wb = openpyxl.load_workbook(self.file_path)
            sheets_data = {}
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                data = [[cell.value for cell in row] for row in ws.iter_rows()]
                sheets_data[sheet_name] = data

            return {
                "content_type": "xlsx",
                "sheets": sheets_data,
                "sheet_names": wb.sheetnames,
            }
        except Exception as e:
            return {
                "content_type": "xlsx",
                "error": str(e),
            }

    async def _process_pptx(self) -> Dict[str, Any]:
        """PowerPoint-Datei"""
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "python-pptx nicht installiert"}

        try:
            prs = Presentation(self.file_path)
            slides_data = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_data.append({
                    "slide": slide_num + 1,
                    "text": "\n".join(slide_text),
                })

            return {
                "content_type": "pptx",
                "slides": slides_data,
                "num_slides": len(prs.slides),
            }
        except Exception as e:
            return {
                "content_type": "pptx",
                "error": str(e),
            }
